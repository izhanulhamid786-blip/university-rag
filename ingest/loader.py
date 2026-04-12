import csv
import hashlib
import json
import logging
import re
import sys
from pathlib import Path

import fitz
from bs4 import BeautifulSoup

sys.path.insert(0, str(Path(__file__).parent.parent))

from rag.settings import get_settings


logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[logging.FileHandler("loader.log", encoding="utf-8"), logging.StreamHandler()],
)
log = logging.getLogger(__name__)


URL_RE = re.compile(r"https?://[^\s'\"<>)]+", re.I)
NOISE_PATTERNS = [
    re.compile(r"original text\s+rate this translation", re.I),
    re.compile(r"your feedback will be used to help improve google translate", re.I),
]
LINK_TAGS = {
    "admission": ["admission", "apply", "enroll", "cuet"],
    "fee": ["fee", "payment", "scholarship"],
    "exam": ["exam", "result", "timetable", "date-sheet", "datesheet"],
    "form": ["form", "download", "bulletin", "prospectus"],
    "notice": ["notice", "circular", "notification"],
}


def _clean_text(text: str) -> str:
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    for pattern in NOISE_PATTERNS:
        text = pattern.sub("", text)

    lines = []
    seen = set()
    for raw_line in text.splitlines():
        line = re.sub(r"\s+", " ", raw_line).strip()
        if not line:
            continue
        key = line.lower()
        if key in seen:
            continue
        seen.add(key)
        lines.append(line)

    return "\n".join(lines).strip()


def _looks_like_error_page(title: str, text: str) -> bool:
    combined = f"{title}\n{text}".lower()
    return any(
        phrase in combined
        for phrase in (
            "server error",
            "404 not found",
            "the resource cannot be found",
            "internal server error",
        )
    )


def _categorize_link(url: str, anchor_text: str) -> list[str]:
    haystack = f"{anchor_text} {url}".lower()
    categories = [name for name, words in LINK_TAGS.items() if any(word in haystack for word in words)]
    return categories or ["general"]


def _unique_links(links: list[dict]) -> list[dict]:
    unique = {}
    for link in links:
        url = (link.get("url") or "").strip()
        if not url:
            continue
        if url not in unique:
            unique[url] = {
                "url": url,
                "anchor_text": (link.get("anchor_text") or url).strip(),
                "source": link.get("source", ""),
                "categories": link.get("categories") or ["general"],
            }
    return list(unique.values())


def _links_from_text(text: str, source: str, soup: BeautifulSoup | None = None) -> list[dict]:
    hrefs = []
    if soup is not None:
        hrefs = [
            (a["href"], a.get_text(" ", strip=True))
            for a in soup.find_all("a", href=True)
            if a["href"].startswith("http")
        ]

    discovered = dict(hrefs + [(url, url) for url in URL_RE.findall(text)])
    return [
        {
            "url": url,
            "anchor_text": anchor_text.strip() or url,
            "source": source,
            "categories": _categorize_link(url, anchor_text),
        }
        for url, anchor_text in discovered.items()
    ]


def _read_text_file(path: Path) -> str:
    for encoding in ("utf-8", "latin-1"):
        try:
            return path.read_text(encoding=encoding)
        except UnicodeDecodeError:
            continue
    return ""


def _structured_links(payload: dict, source: str) -> list[dict]:
    links = []
    for url in payload.get("outlinks", []):
        links.append(
            {
                "url": url,
                "anchor_text": url,
                "source": source,
                "categories": _categorize_link(url, payload.get("title", "")),
            }
        )

    for notice in payload.get("notices", []):
        url = (notice or {}).get("link")
        if not url:
            continue
        anchor_text = (notice or {}).get("text", url)
        links.append(
            {
                "url": url,
                "anchor_text": anchor_text[:200],
                "source": source,
                "categories": _categorize_link(url, anchor_text),
            }
        )

    links.extend(_links_from_text(payload.get("text", ""), source))
    return _unique_links(links)


def _clean_table_cell(cell: object) -> str:
    return re.sub(r"\s+", " ", str(cell or "")).strip()


def _format_table(table: list, index: int) -> str:
    if not isinstance(table, list):
        return ""

    rows = []
    for row in table:
        if not isinstance(row, list):
            continue
        cleaned = [_clean_table_cell(cell) for cell in row]
        if any(cleaned):
            rows.append(cleaned)

    if not rows:
        return ""

    header = rows[0]
    body = rows[1:] if len(rows) > 1 else []
    lines = [f"Table {index}:"]

    if body and any(header):
        for row_number, row in enumerate(body, start=1):
            pairs = []
            for cell_index, value in enumerate(row):
                if not value:
                    continue
                column = header[cell_index] if cell_index < len(header) and header[cell_index] else f"Column {cell_index + 1}"
                pairs.append(f"{column}: {value}")
            if pairs:
                lines.append(f"Row {row_number}: " + " | ".join(pairs))
    else:
        for row_number, row in enumerate(rows, start=1):
            values = [value for value in row if value]
            if values:
                lines.append(f"Row {row_number}: " + " | ".join(values))

    return "\n".join(lines).strip()


def _structured_text(payload: dict) -> str:
    parts = [(payload.get("text") or "").strip()]
    tables = payload.get("tables") or []
    formatted_tables = [_format_table(table, index) for index, table in enumerate(tables, start=1)]
    parts.extend(table for table in formatted_tables if table)
    return "\n\n".join(part for part in parts if part).strip()


def _make_doc(
    *,
    text: str,
    source: str,
    source_path: str,
    source_url: str | None,
    title: str,
    category: str,
    file_type: str,
    links: list[dict],
    source_kind: str,
    scraped_at: str | None = None,
    ocr: bool = False,
    has_table: bool = False,
    table_count: int = 0,
    table_row_count: int = 0,
) -> dict | None:
    clean = _clean_text(text)
    if not clean:
        return None

    canonical_id = source_url or source_path or source
    return {
        "text": clean,
        "links": _unique_links(links),
        "source": source,
        "source_path": source_path,
        "source_url": source_url,
        "title": title.strip() or Path(source_path or source).stem,
        "category": category or "general",
        "file_type": file_type or "txt",
        "doc_id": hashlib.md5(canonical_id.encode("utf-8", errors="ignore")).hexdigest(),
        "source_kind": source_kind,
        "scraped_at": scraped_at,
        "ocr": bool(ocr),
        "has_table": bool(has_table),
        "table_count": int(table_count),
        "table_row_count": int(table_row_count),
    }


def load_structured_records(structured_dir: str | Path | None = None) -> list[dict]:
    settings = get_settings()
    root = Path(structured_dir) if structured_dir else settings.structured_dir
    docs = []
    if not root.exists():
        log.warning("Structured data directory not found: %s", root)
        return docs

    for path in sorted(root.glob("*.json")):
        try:
            payload = json.loads(path.read_text(encoding="utf-8"))
        except Exception as exc:
            log.error("Failed to parse structured record [%s]: %s", path.name, exc)
            continue

        if not isinstance(payload, dict) or "text" not in payload:
            continue

        title = payload.get("title") or path.stem
        source_url = payload.get("url")
        tables = payload.get("tables") or []
        if _looks_like_error_page(title, payload.get("text", "")):
            log.warning("Skipping error page: %s", path.name)
            continue
        doc = _make_doc(
            text=_structured_text(payload),
            source=source_url or str(path),
            source_path=str(path),
            source_url=source_url,
            title=title,
            category=payload.get("category", "general"),
            file_type=payload.get("type", "html"),
            links=_structured_links(payload, source_url or str(path)),
            source_kind="crawler",
            scraped_at=payload.get("scraped_at"),
            ocr=payload.get("ocr", False),
            has_table=bool(tables),
            table_count=len(tables),
            table_row_count=sum(len(table) for table in tables if isinstance(table, list)),
        )
        if doc is None:
            log.warning("Empty structured record: %s", path.name)
            continue

        docs.append(doc)
        log.info("[CRAWLER] %s - %s", path.name, title)

    return docs


def load_pdf(path: Path) -> tuple[str, list[dict]]:
    with fitz.open(path) as doc:
        text = "\n".join(page.get_text() for page in doc)
        links = [
            {
                "url": link["uri"],
                "anchor_text": link["uri"],
                "categories": ["general"],
                "source": str(path),
            }
            for page in doc
            for link in page.get_links()
            if link.get("uri", "").startswith("http")
        ]
    return text, links + _links_from_text(text, str(path))


def load_html(path: Path) -> tuple[str, list[dict]]:
    soup = BeautifulSoup(_read_text_file(path), "html.parser")
    for tag in soup(["nav", "footer", "script", "style"]):
        tag.decompose()
    text = soup.get_text("\n")
    return text, _links_from_text(text, str(path), soup)


def load_txt(path: Path) -> tuple[str, list[dict]]:
    text = _read_text_file(path)
    return text, _links_from_text(text, str(path))


def load_json(path: Path) -> tuple[str, list[dict]]:
    text = json.dumps(json.loads(path.read_text(encoding="utf-8")), indent=2)
    return text, _links_from_text(text, str(path))


def load_csv(path: Path) -> tuple[str, list[dict]]:
    text = "\n".join(
        ", ".join(row) for row in csv.reader(path.read_text(encoding="utf-8", errors="replace").splitlines())
    )
    return text, _links_from_text(text, str(path))


try:
    import docx as _docx

    def load_docx(path: Path) -> tuple[str, list[dict]]:
        doc = _docx.Document(str(path))
        text = "\n".join(paragraph.text for paragraph in doc.paragraphs)
        links = [
            {
                "url": rel._target,
                "anchor_text": rel._target,
                "categories": ["general"],
                "source": str(path),
            }
            for rel in doc.part.rels.values()
            if "hyperlink" in rel.reltype and str(rel._target).startswith("http")
        ]
        return text, links + _links_from_text(text, str(path))

except ImportError:
    def load_docx(path: Path) -> tuple[str, list[dict]]:
        return "", []


try:
    import openpyxl as _openpyxl

    def load_xlsx(path: Path) -> tuple[str, list[dict]]:
        workbook = _openpyxl.load_workbook(str(path), data_only=True)
        text = "\n".join(
            "\t".join(str(cell or "") for cell in row)
            for sheet in workbook
            for row in sheet.iter_rows(values_only=True)
        )
        return text, _links_from_text(text, str(path))

except ImportError:
    def load_xlsx(path: Path) -> tuple[str, list[dict]]:
        return "", []


try:
    from pptx import Presentation as _Presentation

    def load_pptx(path: Path) -> tuple[str, list[dict]]:
        text = "\n".join(
            paragraph.text
            for slide in _Presentation(str(path)).slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            for paragraph in shape.text_frame.paragraphs
        )
        return text, _links_from_text(text, str(path))

except ImportError:
    def load_pptx(path: Path) -> tuple[str, list[dict]]:
        return "", []


RAW_LOADERS = {
    ".pdf": load_pdf,
    ".html": load_html,
    ".htm": load_html,
    ".docx": load_docx,
    ".doc": load_docx,
    ".xlsx": load_xlsx,
    ".xls": load_xlsx,
    ".csv": load_csv,
    ".tsv": load_csv,
    ".txt": load_txt,
    ".md": load_txt,
    ".json": load_json,
    ".pptx": load_pptx,
}


def _raw_doc_from_path(path: Path, source_kind: str) -> dict | None:
    loader = RAW_LOADERS.get(path.suffix.lower())
    if loader is None:
        return None

    text, links = loader(path)
    return _make_doc(
        text=text,
        source=str(path),
        source_path=str(path),
        source_url=None,
        title=path.stem,
        category=path.parent.name,
        file_type=path.suffix.lstrip("."),
        links=links,
        source_kind=source_kind,
    )


def load_manual_files(manual_dir: str | Path | None = None) -> list[dict]:
    settings = get_settings()
    root = Path(manual_dir) if manual_dir else settings.manual_dir
    docs = []
    if not root.exists():
        return docs

    for path in sorted(root.rglob("*")):
        if not path.is_file():
            continue
        doc = _raw_doc_from_path(path, source_kind="manual")
        if doc is None:
            continue
        docs.append(doc)
        log.info("[MANUAL] %s", path.name)
    return docs


def load_legacy_root_files(data_dir: str | Path | None = None) -> list[dict]:
    settings = get_settings()
    root = Path(data_dir) if data_dir else settings.data_dir
    docs = []
    if not root.exists():
        return docs

    for path in sorted(root.iterdir()):
        if not path.is_file():
            continue
        if path.name.startswith("."):
            continue
        doc = _raw_doc_from_path(path, source_kind="legacy_raw")
        if doc is None:
            continue
        docs.append(doc)
        log.info("[LEGACY] %s", path.name)
    return docs


def _dedupe_docs(docs: list[dict]) -> list[dict]:
    chosen = {}
    for doc in docs:
        key = doc.get("source_url") or doc.get("source_path") or doc.get("source")
        existing = chosen.get(key)
        if existing is None:
            chosen[key] = doc
            continue

        current_score = (len(doc.get("text", "")), doc.get("scraped_at") or "")
        existing_score = (len(existing.get("text", "")), existing.get("scraped_at") or "")
        if current_score > existing_score:
            chosen[key] = doc
    return list(chosen.values())


def load_all(data_dir: str | Path | None = None) -> list[dict]:
    settings = get_settings()
    if data_dir:
        custom_root = Path(data_dir)
        structured_dir = custom_root / "structured"
        manual_dir = custom_root / "manual"
        root_dir = custom_root
    else:
        structured_dir = settings.structured_dir
        manual_dir = settings.manual_dir
        root_dir = settings.data_dir

    docs = []
    docs.extend(load_structured_records(structured_dir))

    if settings.include_manual_raw:
        docs.extend(load_manual_files(manual_dir))

    if settings.include_legacy_root_raw:
        docs.extend(load_legacy_root_files(root_dir))

    docs = _dedupe_docs(docs)

    log.info(
        "Done - loaded:%s crawler:%s manual:%s legacy:%s",
        len(docs),
        sum(1 for doc in docs if doc["source_kind"] == "crawler"),
        sum(1 for doc in docs if doc["source_kind"] == "manual"),
        sum(1 for doc in docs if doc["source_kind"] == "legacy_raw"),
    )
    return docs


def find_links(docs: list[dict], query: str) -> list[dict]:
    words = re.findall(r"\w+", query.lower())
    seen = {}
    for doc in docs:
        for link in doc["links"]:
            haystack = f"{link['anchor_text']} {link['url']} {' '.join(link['categories'])}".lower()
            score = sum(1 for word in words if word in haystack)
            if score and (link["url"] not in seen or seen[link["url"]]["_score"] < score):
                seen[link["url"]] = {**link, "_score": score}
    return sorted(seen.values(), key=lambda item: (-item["_score"], item["url"]))


if __name__ == "__main__":
    documents = load_all()
    print(f"\nLoaded: {len(documents)} docs")
    for link in find_links(documents, "admission form cuet ug")[:5]:
        print(f"  {link['anchor_text']} -> {link['url']}")
